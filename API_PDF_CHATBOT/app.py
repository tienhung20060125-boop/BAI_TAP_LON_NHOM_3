import os
import pickle
import tempfile
import uuid

import faiss
import numpy as np
from docx import Document
from dotenv import load_dotenv
from flask import Flask, jsonify, render_template_string, request, session
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pptx import Presentation
from pypdf import PdfReader
from sentence_transformers import SentenceTransformer
from werkzeug.utils import secure_filename


load_dotenv()

app = Flask(__name__)
app.secret_key = os.getenv("FLASK_SECRET_KEY", "local-dev-secret-key")
app.config["MAX_CONTENT_LENGTH"] = 80 * 1024 * 1024

ALLOWED_EXTENSIONS = {".pdf", ".docx", ".pptx"}
CHAT_HISTORIES = {}

llm = ChatGoogleGenerativeAI(model="gemini-2.5-flash")
embedding_model = SentenceTransformer("all-MiniLM-L6-v2")

knowledge_index = None
knowledge_chunks = []


def load_default_knowledge_base():
    if not os.path.exists("faiss_index.bin") or not os.path.exists("chunks.pkl"):
        return None, []

    loaded_index = faiss.read_index("faiss_index.bin")
    with open("chunks.pkl", "rb") as f:
        loaded_chunks = pickle.load(f)

    return loaded_index, loaded_chunks


def get_session_id():
    if "sid" not in session:
        session["sid"] = uuid.uuid4().hex
    CHAT_HISTORIES.setdefault(session["sid"], [])
    return session["sid"]


def get_chat_history():
    return CHAT_HISTORIES.setdefault(get_session_id(), [])


def allowed_file(filename):
    return os.path.splitext(filename.lower())[1] in ALLOWED_EXTENSIONS


def extract_text_from_pdf(path):
    text = ""
    reader = PdfReader(path)
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
    return text


def extract_text_from_docx(path):
    text = ""
    doc = Document(path)
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text


def extract_text_from_pptx(path):
    text = ""
    presentation = Presentation(path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def extract_text_from_uploads(files):
    all_text = ""
    file_names = []

    with tempfile.TemporaryDirectory() as temp_dir:
        for file in files:
            if not file or not file.filename:
                continue

            filename = secure_filename(file.filename)
            if not allowed_file(filename):
                continue

            file_names.append(filename)
            suffix = os.path.splitext(filename.lower())[1]
            temp_path = os.path.join(temp_dir, filename)
            file.save(temp_path)

            if suffix == ".pdf":
                all_text += extract_text_from_pdf(temp_path)
            elif suffix == ".docx":
                all_text += extract_text_from_docx(temp_path)
            elif suffix == ".pptx":
                all_text += extract_text_from_pptx(temp_path)

            all_text += "\n\n"

    return all_text, file_names


def build_knowledge_base(text):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,
        chunk_overlap=100,
    )
    new_chunks = splitter.split_text(text)

    if not new_chunks:
        return None, []

    vectors = embedding_model.encode(new_chunks)
    vectors = np.array(vectors).astype("float32")

    new_index = faiss.IndexFlatL2(vectors.shape[1])
    new_index.add(vectors)

    return new_index, new_chunks


def answer_question(question, chat_history):
    global knowledge_index, knowledge_chunks

    q_vec = embedding_model.encode([question])
    q_vec = np.array(q_vec).astype("float32")

    _, result_indexes = knowledge_index.search(q_vec, k=5)

    unique_idx = []
    for idx in result_indexes[0]:
        idx = int(idx)
        if idx not in unique_idx and 0 <= idx < len(knowledge_chunks):
            unique_idx.append(idx)

    context = ""
    for idx in unique_idx:
        context += knowledge_chunks[idx] + "\n\n"

    history = "\n".join(
        f"User: {q}\nAI: {a}" for q, a in chat_history[-5:]
    )

    prompt = f"""
Bạn là AI Assistant trả lời dựa trên tài liệu được cung cấp.

QUY TẮC:
- Chỉ sử dụng thông tin trong CONTEXT.
- Không được bịa thông tin.
- Nếu không tìm thấy thì trả lời:
  "Không tìm thấy thông tin trong tài liệu."

====================
HISTORY
====================

{history}

====================
CONTEXT
====================

{context}

====================
QUESTION
====================

{question}

====================
ANSWER
====================
"""

    response = llm.invoke(prompt)
    answer = response.content
    sources = [
        {
            "index": idx,
            "text": knowledge_chunks[idx],
        }
        for idx in unique_idx
    ]

    return answer, sources


HTML = """
<!DOCTYPE html>
<html lang="vi">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>AI PDF Assistant</title>
    <style>
        :root {
            --bg: #0f172a;
            --panel: #111827;
            --panel-2: #172033;
            --line: #263244;
            --text: #f8fafc;
            --muted: #94a3b8;
            --blue: #2563eb;
            --blue-dark: #1d4ed8;
            --green: #10b981;
            --amber: #f59e0b;
            --danger: #ef4444;
        }

        * {
            box-sizing: border-box;
            font-family: "Segoe UI", Arial, sans-serif;
            letter-spacing: 0;
        }

        body {
            min-height: 100vh;
            margin: 0;
            background:
                radial-gradient(circle at top left, rgba(37, 99, 235, 0.18), transparent 34rem),
                linear-gradient(180deg, #0f172a 0%, #0b1120 100%);
            color: var(--text);
            display: grid;
            grid-template-columns: 300px 1fr;
        }

        .sidebar {
            min-height: 100vh;
            background: var(--panel);
            border-right: 1px solid var(--line);
            padding: 22px;
            position: sticky;
            top: 0;
        }

        .brand {
            margin-bottom: 24px;
        }

        .brand h2 {
            margin: 0 0 8px;
            font-size: 24px;
        }

        .brand p,
        .muted {
            color: var(--muted);
            line-height: 1.55;
            margin: 0;
        }

        .panel {
            background: rgba(17, 24, 39, 0.84);
            border: 1px solid var(--line);
            border-radius: 8px;
            box-shadow: 0 20px 50px rgba(0, 0, 0, 0.18);
        }

        .upload-box {
            margin-top: 12px;
        }

        .file-input {
            width: 100%;
            min-height: 118px;
            padding: 18px;
            border: 1px dashed #3b4a63;
            border-radius: 8px;
            background: #0f172a;
            color: var(--muted);
            cursor: pointer;
        }

        .file-list {
            display: grid;
            gap: 8px;
            margin: 12px 0;
        }

        .file-chip {
            overflow: hidden;
            text-overflow: ellipsis;
            white-space: nowrap;
            background: #0f172a;
            border: 1px solid var(--line);
            border-radius: 8px;
            padding: 9px 10px;
            color: #dbeafe;
            font-size: 14px;
        }

        button {
            border: 1px solid #334155;
            border-radius: 8px;
            background: #1e293b;
            color: var(--text);
            cursor: pointer;
            font-weight: 700;
            min-height: 44px;
            padding: 0 16px;
            transition: background 160ms ease, border-color 160ms ease, transform 160ms ease;
        }

        button:hover {
            border-color: var(--blue);
            background: var(--blue);
            transform: translateY(-1px);
        }

        button:disabled {
            cursor: not-allowed;
            opacity: 0.55;
            transform: none;
        }

        button.recording {
            background: var(--danger);
            border-color: var(--danger);
        }

        button.recording:hover {
            background: #dc2626;
            border-color: #dc2626;
        }

        .primary {
            width: 100%;
            background: var(--blue);
            border-color: var(--blue);
        }

        .primary:hover {
            background: var(--blue-dark);
        }

        .danger:hover {
            background: var(--danger);
            border-color: var(--danger);
        }

        .stats {
            display: grid;
            gap: 10px;
            margin-top: 22px;
        }

        .stat {
            background: #0f172a;
            border: 1px solid var(--line);
            border-radius: 8px;
            padding: 12px;
        }

        .stat strong {
            display: block;
            font-size: 22px;
        }

        .stat span {
            color: var(--muted);
            font-size: 13px;
        }

        .main {
            min-width: 0;
            padding: 32px;
            display: flex;
            flex-direction: column;
            gap: 16px;
        }

        .hero {
            display: grid;
            grid-template-columns: 1.4fr 0.8fr;
            gap: 16px;
        }

        .hero-main,
        .status-panel,
        .voice-panel,
        .chat-panel {
            padding: 24px;
        }

        .eyebrow {
            color: #93c5fd;
            font-size: 13px;
            font-weight: 800;
            margin-bottom: 10px;
            text-transform: uppercase;
        }

        h1 {
            font-size: clamp(36px, 5vw, 58px);
            line-height: 1.04;
            margin: 0 0 12px;
        }

        h3 {
            margin: 0 0 8px;
        }

        .hero p {
            color: var(--muted);
            font-size: 17px;
            line-height: 1.65;
            margin: 0;
        }

        .status-number {
            font-size: 38px;
            font-weight: 900;
            margin-top: 8px;
        }

        .badge {
            display: inline-flex;
            align-items: center;
            border-radius: 999px;
            font-size: 13px;
            font-weight: 800;
            margin-top: 16px;
            padding: 7px 11px;
        }

        .badge.ready {
            background: rgba(16, 185, 129, 0.14);
            border: 1px solid rgba(16, 185, 129, 0.35);
            color: #86efac;
        }

        .badge.waiting {
            background: rgba(245, 158, 11, 0.14);
            border: 1px solid rgba(245, 158, 11, 0.35);
            color: #fcd34d;
        }

        .voice-panel {
            display: flex;
            align-items: center;
            justify-content: space-between;
            gap: 16px;
        }

        .voice-panel p {
            color: var(--muted);
            margin: 0;
        }

        .voice-actions {
            align-items: center;
            display: flex;
            flex-direction: column;
            gap: 8px;
        }

        .voice-status {
            color: var(--muted);
            font-size: 13px;
            min-height: 18px;
            text-align: center;
        }

        .suggestions {
            display: grid;
            grid-template-columns: repeat(3, 1fr);
            gap: 10px;
        }

        .chat-panel {
            min-height: 420px;
            display: flex;
            flex-direction: column;
        }

        .chat-box {
            display: flex;
            flex: 1;
            flex-direction: column;
            gap: 12px;
            max-height: 52vh;
            overflow-y: auto;
            padding-right: 6px;
        }

        .empty-chat {
            border: 1px dashed #3b4a63;
            border-radius: 8px;
            color: var(--muted);
            line-height: 1.65;
            margin: auto;
            max-width: 620px;
            padding: 28px;
            text-align: center;
        }

        .message {
            border: 1px solid var(--line);
            border-radius: 8px;
            line-height: 1.65;
            max-width: min(760px, 88%);
            padding: 12px 14px;
            white-space: pre-wrap;
        }

        .message.user {
            align-self: flex-end;
            background: var(--blue-dark);
            border-color: var(--blue);
            color: #ffffff;
        }

        .message.bot {
            align-self: flex-start;
            background: #111827;
            color: #e5e7eb;
        }

        .composer {
            border-top: 1px solid var(--line);
            display: flex;
            gap: 10px;
            margin-top: 18px;
            padding-top: 16px;
        }

        .composer input {
            flex: 1;
            min-width: 0;
            border: 1px solid #334155;
            border-radius: 8px;
            background: #ffffff;
            color: #111827;
            font-size: 16px;
            min-height: 46px;
            outline: none;
            padding: 0 14px;
        }

        .sources {
            display: none;
            margin-top: 14px;
        }

        .source-item {
            background: #0f172a;
            border: 1px solid var(--line);
            border-radius: 8px;
            color: #cbd5e1;
            line-height: 1.55;
            margin-top: 8px;
            max-height: 180px;
            overflow: auto;
            padding: 12px;
            white-space: pre-wrap;
        }

        .notice {
            border-radius: 8px;
            display: none;
            line-height: 1.5;
            margin-top: 12px;
            padding: 12px;
        }

        .notice.ok {
            background: rgba(16, 185, 129, 0.12);
            border: 1px solid rgba(16, 185, 129, 0.35);
            color: #bbf7d0;
        }

        .notice.err {
            background: rgba(239, 68, 68, 0.12);
            border: 1px solid rgba(239, 68, 68, 0.35);
            color: #fecaca;
        }

        @media (max-width: 900px) {
            body {
                grid-template-columns: 1fr;
            }

            .sidebar {
                min-height: auto;
                position: static;
            }

            .main {
                padding: 18px;
            }

            .hero,
            .suggestions {
                grid-template-columns: 1fr;
            }

            .voice-panel,
            .composer {
                align-items: stretch;
                flex-direction: column;
            }

            .message {
                max-width: 100%;
            }
        }
    </style>
</head>
<body>
    <aside class="sidebar">
        <div class="brand">
            <h2>AI PDF</h2>
            <p>Tải tài liệu lên, tạo Knowledge Base rồi đặt câu hỏi theo nội dung tài liệu.</p>
        </div>

        <form id="uploadForm" class="upload-box">
            <h3>Tài liệu</h3>
            <input class="file-input" id="files" name="files" type="file" accept=".pdf,.docx,.pptx" multiple>
            <div id="fileList" class="file-list"></div>
            <button class="primary" id="uploadBtn" type="submit">Xử lý tài liệu</button>
            <div id="uploadNotice" class="notice"></div>
        </form>

        <div class="stats">
            <div class="stat">
                <strong id="sidebarChunks">{{ chunk_count }}</strong>
                <span>Chunks hiện có</span>
            </div>
            <div class="stat">
                <strong id="sidebarChats">{{ chat_count }}</strong>
                <span>Lượt hỏi đáp</span>
            </div>
            <button class="danger" id="clearBtn" type="button">Xóa lịch sử chat</button>
        </div>
    </aside>

    <main class="main">
        <section class="hero">
            <div class="hero-main panel">
                <div class="eyebrow">Trợ lý đọc tài liệu</div>
                <h1>AI PDF Assistant</h1>
                <p>
                    Hỏi đáp theo nội dung PDF, DOCX, PPTX bằng tiếng Việt.
                    Tải tài liệu ở thanh bên trái, xử lý Knowledge Base, rồi nhập câu hỏi hoặc dùng giọng nói.
                </p>
            </div>

            <div class="status-panel panel">
                <div class="muted">Knowledge Base</div>
                <div class="status-number" id="heroChunks">{{ chunk_count }}</div>
                    <div class="muted">chunks đang được dùng</div>
                <div id="statusBadge" class="badge {{ 'ready' if kb_ready else 'waiting' }}">
                    {{ 'READY - Sẵn sàng hỏi đáp' if kb_ready else 'WAITING - Cần tải tài liệu' }}
                </div>
            </div>
        </section>

        <section class="voice-panel panel">
            <div>
                <h3>Voice Chat</h3>
                <p>Trình duyệt sẽ nghe giọng nói và đưa câu hỏi vào ô chat.</p>
            </div>
            <div class="voice-actions">
                <button id="voiceBtn" type="button">Bắt đầu nói</button>
                <span id="voiceStatus" class="voice-status">Chưa nghe</span>
            </div>
        </section>

        <section class="suggestions">
            <button class="suggestion" type="button">Tóm tắt tài liệu này</button>
            <button class="suggestion" type="button">Các ý chính quan trọng là gì?</button>
            <button class="suggestion" type="button">Giải thích nội dung dễ hiểu hơn</button>
        </section>

        <section class="chat-panel panel">
            <div id="chatBox" class="chat-box"></div>

            <div id="sources" class="sources">
                <h3>Nội dung được dùng để trả lời</h3>
                <div id="sourceList"></div>
            </div>

            <form id="chatForm" class="composer">
                <input id="question" type="text" placeholder="Nhập câu hỏi về tài liệu..." autocomplete="off">
                <button id="sendBtn" type="submit">Gửi</button>
            </form>
        </section>
    </main>

    <script>
        const initialHistory = {{ chat_history | tojson }};
        let kbReady = {{ 'true' if kb_ready else 'false' }};
        let chatCount = {{ chat_count }};

        const chatBox = document.getElementById("chatBox");
        const chatForm = document.getElementById("chatForm");
        const questionInput = document.getElementById("question");
        const sendBtn = document.getElementById("sendBtn");
        const uploadForm = document.getElementById("uploadForm");
        const filesInput = document.getElementById("files");
        const fileList = document.getElementById("fileList");
        const uploadBtn = document.getElementById("uploadBtn");
        const uploadNotice = document.getElementById("uploadNotice");
        const clearBtn = document.getElementById("clearBtn");
        const voiceBtn = document.getElementById("voiceBtn");
        const voiceStatus = document.getElementById("voiceStatus");
        const sourcePanel = document.getElementById("sources");
        const sourceList = document.getElementById("sourceList");
        let recognition = null;
        let isRecording = false;
        let manualStop = false;
        let voiceTranscript = "";

        function escapeHtml(value) {
            return value
                .replaceAll("&", "&amp;")
                .replaceAll("<", "&lt;")
                .replaceAll(">", "&gt;")
                .replaceAll('"', "&quot;")
                .replaceAll("'", "&#039;");
        }

        function showNotice(element, type, text) {
            element.className = `notice ${type}`;
            element.textContent = text;
            element.style.display = "block";
        }

        function updateStats(chunkCount, nextChatCount, nextKbReady) {
            kbReady = nextKbReady;
            chatCount = nextChatCount;

            document.getElementById("sidebarChunks").textContent = chunkCount;
            document.getElementById("heroChunks").textContent = chunkCount;
            document.getElementById("sidebarChats").textContent = nextChatCount;

            const badge = document.getElementById("statusBadge");
            badge.className = `badge ${kbReady ? "ready" : "waiting"}`;
            badge.textContent = kbReady
                ? "READY - Sẵn sàng hỏi đáp"
                : "WAITING - Cần tải tài liệu";
        }

        function addMessage(role, text) {
            const empty = chatBox.querySelector(".empty-chat");
            if (empty) {
                empty.remove();
            }

            const message = document.createElement("div");
            message.className = `message ${role === "user" ? "user" : "bot"}`;
            message.textContent = text;
            chatBox.appendChild(message);
            chatBox.scrollTop = chatBox.scrollHeight;
        }

        function renderEmptyChat() {
            if (chatBox.children.length > 0) {
                return;
            }

            const empty = document.createElement("div");
            empty.className = "empty-chat";
            empty.textContent = kbReady
                ? "Đặt câu hỏi đầu tiên để bắt đầu hỏi đáp theo tài liệu."
                : "Hãy tải tài liệu và bấm Xử lý tài liệu trước khi đặt câu hỏi.";
            chatBox.appendChild(empty);
        }

        function renderHistory() {
            chatBox.innerHTML = "";
            initialHistory.forEach(([question, answer]) => {
                addMessage("user", question);
                addMessage("bot", answer);
            });
            renderEmptyChat();
        }

        function renderSources(sources) {
            sourceList.innerHTML = "";
            if (!sources || sources.length === 0) {
                sourcePanel.style.display = "none";
                return;
            }

            sources.forEach((source) => {
                const item = document.createElement("div");
                item.className = "source-item";
                item.textContent = `Chunk ${source.index}\\n\\n${source.text}`;
                sourceList.appendChild(item);
            });
            sourcePanel.style.display = "block";
        }

        filesInput.addEventListener("change", () => {
            fileList.innerHTML = "";
            Array.from(filesInput.files).forEach((file) => {
                const item = document.createElement("div");
                item.className = "file-chip";
                item.textContent = `FILE - ${file.name}`;
                fileList.appendChild(item);
            });
        });

        uploadForm.addEventListener("submit", async (event) => {
            event.preventDefault();

            if (!filesInput.files.length) {
                showNotice(uploadNotice, "err", "Hãy chọn ít nhất một file.");
                return;
            }

            uploadBtn.disabled = true;
            uploadBtn.textContent = "Đang xử lý...";

            const formData = new FormData();
            Array.from(filesInput.files).forEach((file) => {
                formData.append("files", file);
            });

            try {
                const response = await fetch("/upload", {
                    method: "POST",
                    body: formData,
                });
                const data = await response.json();

                if (!response.ok) {
                    throw new Error(data.error || "Không xử lý được tài liệu.");
                }

                updateStats(data.chunk_count, chatCount, true);
                showNotice(uploadNotice, "ok", data.message);
                renderEmptyChat();
            } catch (error) {
                showNotice(uploadNotice, "err", error.message);
            } finally {
                uploadBtn.disabled = false;
                uploadBtn.textContent = "Xử lý tài liệu";
            }
        });

        chatForm.addEventListener("submit", async (event) => {
            event.preventDefault();

            const question = questionInput.value.trim();
            if (!question) {
                return;
            }

            if (!kbReady) {
                addMessage("bot", "Chưa có Knowledge Base. Hãy tải tài liệu lên và bấm Xử lý tài liệu.");
                return;
            }

            addMessage("user", question);
            questionInput.value = "";
            sendBtn.disabled = true;
            sendBtn.textContent = "Đang trả lời...";
            sourcePanel.style.display = "none";

            try {
                const response = await fetch("/ask", {
                    method: "POST",
                    headers: {
                        "Content-Type": "application/json",
                    },
                    body: JSON.stringify({ question }),
                });
                const data = await response.json();

                if (!response.ok) {
                    throw new Error(data.error || "Không tạo được câu trả lời.");
                }

                addMessage("bot", data.answer);
                renderSources(data.sources);
                updateStats(data.chunk_count, data.chat_count, data.kb_ready);
            } catch (error) {
                addMessage("bot", error.message);
            } finally {
                sendBtn.disabled = false;
                sendBtn.textContent = "Gửi";
                questionInput.focus();
            }
        });

        document.querySelectorAll(".suggestion").forEach((button) => {
            button.addEventListener("click", () => {
                questionInput.value = button.textContent.trim();
                questionInput.focus();
            });
        });

        clearBtn.addEventListener("click", async () => {
            await fetch("/clear", { method: "POST" });
            chatBox.innerHTML = "";
            sourcePanel.style.display = "none";
            updateStats(document.getElementById("heroChunks").textContent, 0, kbReady);
            renderEmptyChat();
        });

        function finishVoiceRecording() {
            isRecording = false;
            manualStop = true;
            voiceBtn.disabled = false;
            voiceBtn.classList.remove("recording");
            voiceBtn.textContent = "Bắt đầu nói";
            voiceStatus.textContent = "Đã dừng nghe";

            if (recognition) {
                recognition.stop();
            }

            if (voiceTranscript.trim()) {
                questionInput.value = voiceTranscript.trim();
                questionInput.focus();
            }
        }

        function startVoiceRecording() {
            const SpeechRecognition = window.SpeechRecognition || window.webkitSpeechRecognition;
            if (!SpeechRecognition) {
                addMessage("bot", "Trình duyệt này không hỗ trợ Web Speech API. Hãy nhập câu hỏi bằng bàn phím.");
                return;
            }

            recognition = new SpeechRecognition();
            recognition.lang = "vi-VN";
            recognition.continuous = true;
            recognition.interimResults = true;
            recognition.maxAlternatives = 1;

            isRecording = true;
            manualStop = false;
            voiceTranscript = questionInput.value.trim();
            voiceBtn.disabled = false;
            voiceBtn.classList.add("recording");
            voiceBtn.textContent = "Dừng nghe";
            voiceStatus.textContent = "Đang nghe... bấm Dừng nghe để kết thúc";

            recognition.onresult = (event) => {
                let finalText = "";
                let interimText = "";

                for (let i = event.resultIndex; i < event.results.length; i += 1) {
                    const transcript = event.results[i][0].transcript.trim();
                    if (event.results[i].isFinal) {
                        finalText += ` ${transcript}`;
                    } else {
                        interimText += ` ${transcript}`;
                    }
                }

                if (finalText.trim()) {
                    voiceTranscript = `${voiceTranscript} ${finalText}`.trim();
                }

                questionInput.value = `${voiceTranscript} ${interimText}`.trim();
            };

            recognition.onerror = (event) => {
                if (manualStop) {
                    return;
                }

                if (event.error === "no-speech" && isRecording && !manualStop) {
                    return;
                }

                isRecording = false;
                manualStop = true;
                voiceBtn.disabled = false;
                voiceBtn.classList.remove("recording");
                voiceBtn.textContent = "Bắt đầu nói";
                voiceStatus.textContent = "Đã dừng nghe";
                addMessage("bot", "Không nhận diện được giọng nói. Hãy thử lại hoặc nhập bằng bàn phím.");
            };

            recognition.onend = () => {
                if (isRecording && !manualStop) {
                    try {
                        recognition.start();
                        voiceBtn.classList.add("recording");
                        voiceBtn.textContent = "Dừng nghe";
                        voiceStatus.textContent = "Đang nghe... bấm Dừng nghe để kết thúc";
                    } catch (error) {
                        isRecording = false;
                        voiceBtn.disabled = false;
                        voiceBtn.classList.remove("recording");
                        voiceBtn.textContent = "Bắt đầu nói";
                        voiceStatus.textContent = "Đã dừng nghe";
                    }
                    return;
                }

                voiceBtn.disabled = false;
                voiceBtn.classList.remove("recording");
                voiceBtn.textContent = "Bắt đầu nói";
                voiceStatus.textContent = "Đã dừng nghe";
            };

            try {
                recognition.start();
            } catch (error) {
                isRecording = false;
                manualStop = true;
                voiceBtn.disabled = false;
                voiceBtn.classList.remove("recording");
                voiceBtn.textContent = "Bắt đầu nói";
                voiceStatus.textContent = "Không bật được microphone";
                addMessage("bot", "Không thể bật ghi âm. Hãy kiểm tra quyền microphone của trình duyệt.");
            }
        }

        voiceBtn.addEventListener("click", () => {
            if (isRecording) {
                finishVoiceRecording();
            } else {
                startVoiceRecording();
            }
        });

        renderHistory();
    </script>
</body>
</html>
"""


@app.route("/")
def home():
    chat_history = get_chat_history()
    kb_ready = knowledge_index is not None and len(knowledge_chunks) > 0

    return render_template_string(
        HTML,
        chat_history=chat_history,
        chat_count=len(chat_history),
        chunk_count=len(knowledge_chunks),
        kb_ready=kb_ready,
    )


@app.post("/upload")
def upload():
    global knowledge_index, knowledge_chunks

    files = request.files.getlist("files")
    if not files:
        return jsonify({"error": "Không có file nào được tải lên."}), 400

    invalid_files = [
        secure_filename(file.filename)
        for file in files
        if file and file.filename and not allowed_file(file.filename)
    ]
    if invalid_files:
        return jsonify({"error": "Chỉ hỗ trợ PDF, DOCX và PPTX."}), 400

    try:
        all_text, file_names = extract_text_from_uploads(files)
        if not file_names:
            return jsonify({"error": "Hãy chọn file PDF, DOCX hoặc PPTX hợp lệ."}), 400

        if not all_text.strip():
            return jsonify({"error": "Không đọc được nội dung từ tài liệu."}), 400

        new_index, new_chunks = build_knowledge_base(all_text)
        if new_index is None:
            return jsonify({"error": "Không tạo được Knowledge Base."}), 400

        knowledge_index = new_index
        knowledge_chunks = new_chunks

        return jsonify(
            {
                "message": f"Đã xử lý {len(file_names)} file và tạo {len(new_chunks)} chunks.",
                "chunk_count": len(knowledge_chunks),
                "kb_ready": True,
            }
        )
    except Exception as exc:
        return jsonify({"error": f"Lỗi xử lý tài liệu: {exc}"}), 500


@app.post("/ask")
def ask():
    kb_ready = knowledge_index is not None and len(knowledge_chunks) > 0
    if not kb_ready:
        return jsonify({"error": "Chưa có Knowledge Base. Hãy tải tài liệu trước."}), 400

    data = request.get_json(silent=True) or {}
    question = (data.get("question") or "").strip()
    if not question:
        return jsonify({"error": "Câu hỏi đang trống."}), 400

    chat_history = get_chat_history()

    try:
        answer, sources = answer_question(question, chat_history)
    except Exception as exc:
        return jsonify({"error": f"Lỗi tạo câu trả lời: {exc}"}), 500

    chat_history.append((question, answer))

    return jsonify(
        {
            "answer": answer,
            "sources": sources,
            "chat_count": len(chat_history),
            "chunk_count": len(knowledge_chunks),
            "kb_ready": True,
        }
    )


@app.post("/clear")
def clear():
    CHAT_HISTORIES[get_session_id()] = []
    return jsonify({"ok": True})


if __name__ == "__main__":
    knowledge_index, knowledge_chunks = load_default_knowledge_base()
    app.run(host="0.0.0.0", port=5000, debug=True)